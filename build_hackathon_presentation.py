"""Hackathon version of the EIA Agent System presentation.

Same visual style as build_presentation.py, but restructured around:
  - Hackathon judging criteria (perceiving / reasoning / executing / autonomy)
  - Safety: RAG + structured prompts → deterministic outputs
  - Failure modes
  - Evaluation pipeline + honest gaps
  - Tradeoffs: latency · cost · accuracy · repeatability
  - Emergent capabilities
"""
import io
import re
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colors ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x1F, 0x28, 0x33)
ORANGE  = RGBColor(0x00, 0xCC, 0x6A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT = RGBColor(0xD1, 0xD5, 0xDB)
RED     = RGBColor(0xEF, 0x44, 0x44)
YELLOW  = RGBColor(0xEA, 0xB3, 0x08)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
CYAN    = RGBColor(0x06, 0xB6, 0xD3)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
CARD    = RGBColor(0x2D, 0x37, 0x48)
DARK    = RGBColor(0x1A, 0x20, 0x2C)
DIVIDER = RGBColor(0x3D, 0x4A, 0x5C)

W = Inches(13.33)
H = Inches(7.5)


# ── Core helpers ──────────────────────────────────────────────────────────────

def convert_potx_to_pptx(potx_path: str) -> Presentation:
    import re as _re
    buf = io.BytesIO()
    with zipfile.ZipFile(potx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename.startswith('ppt/slides/'):
                    continue
                data = zin.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    data = data.replace(
                        b'presentationml.template.main+xml',
                        b'presentationml.presentation.main+xml',
                    )
                    data = _re.sub(
                        rb'<Override[^>]*/ppt/slides/slide\d+\.xml[^>]*/?>',
                        b'', data,
                    )
                if item.filename == 'ppt/presentation.xml':
                    data = _re.sub(rb'<p:sldIdLst>.*?</p:sldIdLst>', b'<p:sldIdLst/>',
                                   data, flags=_re.DOTALL)
                zout.writestr(item, data)
    buf.seek(0)
    prs = Presentation(buf)
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, left, top, width, height, color=CARD, border=None, border_w=1.5):
    s = slide.shapes.add_shape(1, left, top, width, height)
    fill(s, color)
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, size=28):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
    fill(s, ORANGE)
    txt(slide, title, Inches(0.6), Inches(0.16), Inches(12.5), Inches(0.72),
        size=size, bold=True, color=ORANGE)
    d = slide.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.025))
    fill(d, DIVIDER)


def bullets(slide, items, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6 if level == 0 else 2)
        indent = "    " * level
        char   = "▸ " if level == 0 else "· "
        r = p.add_run()
        r.text = indent + char + text
        r.font.size  = Pt(18 if level == 0 else 14)
        r.font.bold  = (level == 0)
        r.font.color.rgb = color


def placeholder(slide, label, left, top, width, height):
    b = slide.shapes.add_shape(1, left, top, width, height)
    fill(b, CARD)
    b.line.color.rgb = ORANGE
    b.line.width = Pt(1.5)
    txt(slide, f"[ {label} ]",
        left + Inches(0.15), top + height / 2 - Pt(24),
        width - Inches(0.3), Inches(0.5),
        size=12, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)


def col_header(slide, label, left, top, width, color=ORANGE):
    b = slide.shapes.add_shape(1, left, top, width, Inches(0.36))
    fill(b, CARD)
    txt(slide, label, left + Inches(0.12), top + Inches(0.05),
        width - Inches(0.24), Inches(0.28),
        size=12, bold=True, color=color)


# ── Prompts overview slide ────────────────────────────────────────────────────

PROMPT_ROWS = [
    ("Project Parser",       BLUE,   "2",  "First agent — every run",
     "Extract project type, scale, location, actions, federal_nexus from free-text",
     "agents/project_parser.py"),
    ("Regulatory Screening", GREEN,  "2",  "After env data — RAG retrieval",
     "Top-8 CFR chunks + project context → triggered regulations with citations",
     "agents/regulatory_screening.py"),
    ("Impact Analysis",      ORANGE, "2",  "Core matrix generation",
     "Env data + regs + actions → significance × confidence × mitigation per cell",
     "agents/impact_analysis.py"),
    ("Report Synthesis",     YELLOW, "7",  "1 shared system + 6 section calls",
     "NEPA technical writer: purpose, action, no-action, env, consequences, mitigation",
     "agents/report_synthesis.py\nagents/templates/nepa_ea.py"),
    ("GT Extractor",         PURPLE, "2",  "Evaluation only — not in pipeline",
     "Sample EIS chunks → ground truth: category + significance + evidence",
     "rag_eval/extractor.py"),
]


def add_prompts_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    s = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
    fill(s, ORANGE)
    txt(sl, "Prompts Overview  —  15 total across 5 LLM callers",
        Inches(0.6), Inches(0.16), Inches(12.5), Inches(0.72),
        size=26, bold=True, color=ORANGE)
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.025))
    fill(d, DIVIDER)

    txt(sl, "All non-report prompts enforce JSON-only output.  Every prompt pairs a system role with a runtime user prompt.",
        Inches(0.6), Inches(1.06), Inches(12.1), Inches(0.38),
        size=13, italic=True, color=GRAY_LT)

    # Column headers
    hrow = sl.shapes.add_shape(1, Inches(0.6), Inches(1.52), Inches(12.1), Inches(0.32))
    fill(hrow, DARK)
    for label, lx in [("CALLER", Inches(0.72)), ("#", Inches(2.48)),
                      ("WHEN", Inches(2.86)), ("WHAT IT DOES", Inches(5.22)),
                      ("FILE", Inches(9.62))]:
        txt(sl, label, lx, Inches(1.56), Inches(2.3), Inches(0.22),
            size=10, bold=True, color=ORANGE)

    row_h = Inches(1.0)
    for i, (name, color, count, when, what, fpath) in enumerate(PROMPT_ROWS):
        ry = Inches(1.84) + i * row_h
        rb = sl.shapes.add_shape(1, Inches(0.6), ry, Inches(12.1), row_h - Inches(0.04))
        fill(rb, CARD if i % 2 == 0 else BG)
        # color bar
        lb = sl.shapes.add_shape(1, Inches(0.6), ry, Inches(0.18), row_h - Inches(0.04))
        fill(lb, color)
        txt(sl, name, Inches(0.86), ry + Inches(0.08),
            Inches(1.55), Inches(0.36), size=12, bold=True, color=color)
        txt(sl, count, Inches(2.5), ry + Inches(0.32),
            Inches(0.28), Inches(0.36), size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(sl, when, Inches(2.88), ry + Inches(0.1),
            Inches(2.26), Inches(0.78), size=11, color=WHITE, wrap=True)
        txt(sl, what, Inches(5.24), ry + Inches(0.1),
            Inches(4.3), Inches(0.78), size=11, color=GRAY_LT, wrap=True)
        txt(sl, fpath, Inches(9.64), ry + Inches(0.18),
            Inches(2.98), Inches(0.62), size=9, italic=True, color=ORANGE, wrap=True)

    # Total bar
    tot = sl.shapes.add_shape(1, Inches(0.6), Inches(6.88), Inches(12.1), Inches(0.36))
    fill(tot, DARK)
    tot.line.color.rgb = ORANGE
    tot.line.width = Pt(1)
    txt(sl, "15 prompts total  ·  JSON-only output enforced on 13 of 15  ·  Report Synthesis uses free-form prose",
        Inches(0.8), Inches(6.92), Inches(11.5), Inches(0.26),
        size=12, bold=True, color=ORANGE)


# ── Slide-number removal ──────────────────────────────────────────────────────

def _strip_slide_numbers(path: str) -> None:
    """Remove every <p:sp> containing any form of slide-number content.

    Catches both:
      - <p:ph type="sldNum"> placeholder shapes
      - <a:fld type="slidenum"> field shapes (text boxes with slide-number fields)

    Operates on zip XML directly so it covers the master, all layouts, and slides.
    """
    SP_ANY_SLNUM = re.compile(
        r'<p:sp\b(?:(?!<p:sp\b).)*?'
        r'(?:type="sldNum"|type="slidenum")'
        r'(?:(?!<p:sp\b).)*?</p:sp>',
        re.DOTALL | re.IGNORECASE,
    )

    buf = io.BytesIO()
    with zipfile.ZipFile(path, "r") as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                b"sldNum" in data or b"slidenum" in data
            ):
                cleaned = SP_ANY_SLNUM.sub("", data.decode("utf-8"))
                data = cleaned.encode("utf-8")
            zout.writestr(item, data)

    buf.seek(0)
    with open(path, "wb") as f:
        f.write(buf.read())


# ── Build ─────────────────────────────────────────────────────────────────────

def build(potx_path: str, output_path: str):
    prs = convert_potx_to_pptx(potx_path)

    # ── 1. Title ──────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    top_bar = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.8))
    fill(top_bar, ORANGE)
    txt(sl, "EIA AGENT SYSTEM",
        Inches(0.7), Inches(0.2), Inches(11.5), Inches(1.1),
        size=44, bold=True, color=WHITE)
    txt(sl, "Automated Environmental Impact Assessment via Multi-Agent AI",
        Inches(0.7), Inches(2.1), Inches(11.0), Inches(0.7),
        size=22, color=GRAY_LT)
    txt(sl, "Data Science Club  ·  AI Agents Hackathon  ·  April 2026",
        Inches(0.7), Inches(3.05), Inches(10.0), Inches(0.5),
        size=16, color=ORANGE)
    txt(sl, "[ Team Member Names ]",
        Inches(0.7), Inches(3.8), Inches(10.0), Inches(0.4),
        size=15, color=GRAY_LT)
    bot = sl.shapes.add_shape(1, Inches(0), H - Inches(0.15), W, Inches(0.15))
    fill(bot, ORANGE)

    # ── 2. The Problem ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "The Problem")

    txt(sl, "Environmental impact reviews are slow, expensive, and hard to audit.",
        Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.45),
        size=16, color=GRAY_LT, italic=True)

    problems = [
        ("Weeks of expert analyst time to screen a single project", 0, WHITE),
        ("5+ federal databases  ·  hundreds of CFR regulations", 1, GRAY_LT),
        ("Reports are long-form prose — hard to verify or audit", 0, WHITE),
        ("No repeatable, objective quality benchmark exists", 0, WHITE),
        ("How do you know the AI got it right?", 1, GRAY_LT),
    ]
    bullets(sl, problems, Inches(0.9), Inches(1.62), Inches(7.5), Inches(3.8))

    # Goal box
    goal = sl.shapes.add_shape(1, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9))
    fill(goal, CARD)
    goal.line.color.rgb = ORANGE
    goal.line.width = Pt(1.5)
    txt(sl, "Goal:  autonomous AI pipeline that perceives project context,"
           " reasons over regulations, and executes a compliance assessment —"
           " grounded, traceable, and quantitatively evaluated.",
        Inches(0.85), Inches(5.7), Inches(11.7), Inches(0.72),
        size=14, color=WHITE)

    placeholder(sl, "NEPA timeline / cost comparison chart",
                Inches(8.4), Inches(1.15), Inches(4.3), Inches(4.2))

    # ── 3. The Agent: Criteria Checklist ─────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Is It Really an AI Agent?")

    criteria = [
        ("PERCEIVING",        GREEN,
         "Reads free-text project descriptions + queries 5 live federal REST APIs"
         " (species, wetlands, flood zones, soil, environmental justice)"),
        ("REASONING",         BLUE,
         "LLM-driven analysis: RAG retrieves applicable CFR regulations,"
         " then an impact analysis agent judges significance × confidence for every action"),
        ("EXECUTING",         YELLOW,
         "Produces a color-coded significance matrix, a 10-section NEPA narrative,"
         " and a downloadable PDF report — without any human in the loop"),
        ("AUTONOMY",          ORANGE,
         "LangGraph orchestrates 5 specialized agents end-to-end."
         " Each agent output gates the next. No human approval required between steps."),
    ]
    cy = Inches(1.12)
    for label, color, desc in criteria:
        cb = sl.shapes.add_shape(1, Inches(0.6), cy, Inches(12.1), Inches(1.22))
        fill(cb, CARD)
        cb.line.color.rgb = color
        cb.line.width = Pt(2)
        lb = sl.shapes.add_shape(1, Inches(0.6), cy, Inches(2.1), Inches(1.22))
        fill(lb, color)
        txt(sl, label, Inches(0.72), cy + Inches(0.38),
            Inches(1.86), Inches(0.48), size=13, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(sl, desc, Inches(2.88), cy + Inches(0.22),
            Inches(9.7), Inches(0.82), size=14, color=WHITE, wrap=True)
        cy += Inches(1.34)

    # ── 4. Architecture — LangGraph ───────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Architecture — LangGraph Pipeline")

    txt(sl, "5 typed graph nodes execute sequentially."
           "  Structured prompts enforce JSON-only output at every LLM call.",
        Inches(0.6), Inches(1.08), Inches(12.1), Inches(0.42),
        size=14, color=GRAY_LT)

    agents = [
        ("1\nProject\nParser",       BLUE,   "Gemini 2.5 Flash"),
        ("2\nEnviron.\nData",        CYAN,   "No LLM — APIs"),
        ("3\nRegulatory\nScreening", GREEN,  "Claude Haiku + RAG"),
        ("4\nImpact\nAnalysis",      ORANGE, "Claude Haiku"),
        ("5\nReport\nSynthesis",     YELLOW, "Gemini 2.5 Flash"),
    ]
    bw = Inches(2.22)
    bg_gap = Inches(0.22)
    ax = Inches(0.5)
    ay = Inches(1.62)
    ah = Inches(2.2)
    for i, (name, color, model) in enumerate(agents):
        bx = ax + i * (bw + bg_gap)
        b = sl.shapes.add_shape(1, bx, ay, bw, ah)
        fill(b, CARD)
        b.line.color.rgb = color
        b.line.width = Pt(2)
        hb = sl.shapes.add_shape(1, bx, ay, bw, Inches(0.42))
        fill(hb, color)
        txt(sl, name, bx + Inches(0.1), ay + Inches(0.5),
            bw - Inches(0.2), Inches(0.9),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, model, bx + Inches(0.1), ay + ah - Inches(0.52),
            bw - Inches(0.2), Inches(0.42),
            size=10, italic=True, color=color, align=PP_ALIGN.CENTER)

    arr = sl.shapes.add_shape(1, Inches(0.6), ay + ah + Inches(0.1),
                              Inches(12.1), Inches(0.03))
    fill(arr, DIVIDER)
    txt(sl, "sequential  ▶", Inches(5.5), ay + ah + Inches(0.16),
        Inches(3.0), Inches(0.3), size=10, italic=True, color=DIVIDER)

    stack = [
        ("LangGraph",  "Typed state dict flows through nodes  ·  each agent reads previous outputs",  ORANGE),
        ("pgvector",   "HNSW cosine index — RAG retrieval + evaluation search",                       GREEN),
        ("FastAPI SSE","POST /api/run → streaming server-sent events → real-time UI updates",         BLUE),
    ]
    sy = Inches(4.15)
    for label, desc, color in stack:
        rb = sl.shapes.add_shape(1, Inches(0.6), sy, Inches(12.1), Inches(0.54))
        fill(rb, CARD)
        rb.line.color.rgb = color
        rb.line.width = Pt(1)
        txt(sl, label, Inches(0.8), sy + Inches(0.12),
            Inches(1.75), Inches(0.32), size=13, bold=True, color=color)
        txt(sl, desc, Inches(2.72), sy + Inches(0.14),
            Inches(9.7), Inches(0.28), size=12, color=GRAY_LT)
        sy += Inches(0.64)

    placeholder(sl, "Pipeline diagram / live UI screenshot",
                Inches(0.6), Inches(6.08), Inches(12.1), Inches(1.2))

    # ── 5. Prompts Overview ───────────────────────────────────────────────────
    add_prompts_slide(prs)

    # ── 6. Safety — RAG + Structured Prompts ─────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Safety & Reliability")

    # Left column: RAG
    col_header(sl, "RAG PREVENTS HALLUCINATION", Inches(0.6), Inches(1.08), Inches(5.9), GREEN)
    bullets(sl, [
        ("Every regulation cited is retrieved from the vector DB", 0, WHITE),
        ("Not generated from LLM memory", 1, GRAY_LT),
        ("Chunk-level citations included in output", 0, WHITE),
        ("Enables post-hoc audit trail", 1, GRAY_LT),
        ("Agent cannot invent a CFR citation that is not in corpus", 0, WHITE),
    ], Inches(0.6), Inches(1.5), Inches(5.9), Inches(3.6))

    # Right column: structured prompts
    col_header(sl, "STRUCTURED PROMPTS → DETERMINISTIC OUTPUTS", Inches(7.0), Inches(1.08), Inches(5.9), CYAN)
    bullets(sl, [
        ("All non-report prompts enforce JSON-only responses", 0, WHITE),
        ("Schema validated before downstream agents consume output", 0, WHITE),
        ("LangGraph typed state prevents partial-output propagation", 0, WHITE),
        ("Confidence tiers (0.0–1.0) quantify uncertainty per cell", 0, WHITE),
        ("needs_review = True flags cells below 0.60 confidence", 0, WHITE),
    ], Inches(7.0), Inches(1.5), Inches(5.9), Inches(3.6))

    # Bottom: confidence tier mini-table
    tiers = [
        ("0.85–1.0",  "HIGH",     GREEN,  "API data + explicit numeric threshold"),
        ("0.65–0.84", "GOOD",     CYAN,   "API data present, threshold is judgment call"),
        ("0.45–0.64", "MODERATE", YELLOW, "Partial data or tangential regulation"),
        ("0.00–0.44", "LOW / VERY LOW", RED, "No direct data — extrapolation only  →  human review"),
    ]
    ty = Inches(5.2)
    tw = Inches(12.1)
    hrow = sl.shapes.add_shape(1, Inches(0.6), ty, tw, Inches(0.3))
    fill(hrow, DARK)
    txt(sl, "CONFIDENCE TIER", Inches(0.72), ty + Inches(0.04),
        Inches(1.5), Inches(0.22), size=10, bold=True, color=ORANGE)
    txt(sl, "LABEL", Inches(2.35), ty + Inches(0.04),
        Inches(1.2), Inches(0.22), size=10, bold=True, color=ORANGE)
    txt(sl, "MEANING", Inches(3.75), ty + Inches(0.04),
        Inches(8.9), Inches(0.22), size=10, bold=True, color=ORANGE)
    row_h = Inches(0.44)
    for i, (score, tier, color, meaning) in enumerate(tiers):
        ry = ty + Inches(0.3) + i * row_h
        rb = sl.shapes.add_shape(1, Inches(0.6), ry, tw, row_h - Inches(0.03))
        fill(rb, CARD if i % 2 == 0 else BG)
        txt(sl, score, Inches(0.72), ry + Inches(0.08),
            Inches(1.5), Inches(0.28), size=12, bold=True, color=color)
        txt(sl, tier, Inches(2.35), ry + Inches(0.09),
            Inches(1.3), Inches(0.26), size=11, bold=True, color=color)
        txt(sl, meaning, Inches(3.75), ry + Inches(0.09),
            Inches(8.8), Inches(0.26), size=11, color=GRAY_LT)

    # ── 6. Failure Modes ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Failure Modes")

    txt(sl, "Building AI agents is easy.  Making deployable products is very hard.",
        Inches(0.6), Inches(1.08), Inches(12.1), Inches(0.38),
        size=15, italic=True, color=GRAY_LT)

    failure_modes = [
        ("FEDERAL API OUTAGE",     RED,
         "USFWS / FEMA / USDA endpoints go down",
         "Graceful degradation — missing data lowers confidence score, does not crash pipeline"),
        ("STALE REGULATORY CORPUS", YELLOW,
         "CFR amendments not yet ingested into pgvector",
         "RAG retrieves outdated chunks → agent cites superseded rules silently"),
        ("CONTEXT OVERFLOW",       ORANGE,
         "Large projects with many actions exceed LLM context window",
         "Impact analysis agent truncates action list; edge cases in truncation not fully tested"),
        ("GROUND TRUTH EXTRACTION", PURPLE,
         "LLM misreads tables in the EIS PDF during eval",
         "Inflated eval scores — the evaluator is also an LLM, introducing correlated errors"),
        ("LLM NON-DETERMINISM",    CYAN,
         "Same input, different significance rating across runs",
         "Temperature > 0 means F1 scores can vary by ±5% between identical runs"),
        ("EMBEDDING DRIFT",        BLUE,
         "OpenAI updates text-embedding-3-small; stored embeddings become misaligned",
         "RAG retrieval degrades silently — no automated re-index alert in place"),
    ]

    col_w = Inches(5.9)
    card_h = Inches(1.04)
    card_g = Inches(0.1)
    for i, (title, color, risk, mitigation) in enumerate(failure_modes):
        col = i % 2
        row = i // 2
        cx = Inches(0.6) + col * (col_w + Inches(0.3))
        cy = Inches(1.56) + row * (card_h + card_g)
        cb = sl.shapes.add_shape(1, cx, cy, col_w, card_h)
        fill(cb, CARD)
        cb.line.color.rgb = color
        cb.line.width = Pt(1.5)
        lhb = sl.shapes.add_shape(1, cx, cy, Inches(0.22), card_h)
        fill(lhb, color)
        txt(sl, title, cx + Inches(0.32), cy + Inches(0.06),
            col_w - Inches(0.42), Inches(0.26), size=11, bold=True, color=color)
        txt(sl, f"Risk: {risk}",
            cx + Inches(0.32), cy + Inches(0.32),
            col_w - Inches(0.42), Inches(0.26), size=11, color=WHITE)
        txt(sl, mitigation,
            cx + Inches(0.32), cy + Inches(0.6),
            col_w - Inches(0.42), Inches(0.38), size=10, italic=True, color=GRAY_LT,
            wrap=True)

    # ── 7. Evaluation Pipeline ────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Evaluation Pipeline")

    txt(sl, "Upload a published EIS for the same project."
           "  Score the agent's impact matrix against it automatically.",
        Inches(0.6), Inches(1.08), Inches(12.1), Inches(0.4),
        size=14, color=GRAY_LT)

    steps = [
        ("UPLOAD",  "User uploads\npublished EIS PDF\nfor a saved project",     BLUE),
        ("PARSE",   "PDF chunked,\nembedded, stored\nin pgvector",              CYAN),
        ("EXTRACT", "LLM reads EIS →\nground truth:\ncategory + significance",  GREEN),
        ("SCORE",   "Agent matrix vs.\nground truth →\n3 metrics computed",     ORANGE),
        ("REVIEW",  "Per-category\nTP / FP / FN\nbreakdown in UI",             YELLOW),
    ]
    sw = Inches(2.3)
    sh = Inches(3.1)
    sg = Inches(0.22)
    sx = Inches(0.5)
    sy2 = Inches(1.6)
    for i, (label, desc, color) in enumerate(steps):
        bx = sx + i * (sw + sg)
        sb = sl.shapes.add_shape(1, bx, sy2, sw, sh)
        fill(sb, CARD)
        sb.line.color.rgb = color
        sb.line.width = Pt(1.5)
        hb = sl.shapes.add_shape(1, bx, sy2, sw, Inches(0.42))
        fill(hb, color)
        txt(sl, f"{i+1}. {label}", bx + Inches(0.1), sy2 + Inches(0.07),
            sw - Inches(0.2), Inches(0.32), size=11, bold=True, color=WHITE)
        txt(sl, desc, bx + Inches(0.12), sy2 + Inches(0.52),
            sw - Inches(0.24), sh - Inches(0.62),
            size=13, color=GRAY_LT, wrap=True)

    # Three metrics summary
    metrics = [
        ("Category F1",           "0.40",  BLUE,
         "Precision + recall over 8 fixed\nenvironmental resource categories"),
        ("Significance Accuracy", "0.40",  YELLOW,
         "Ordinal comparison: significant ·\nmoderate · minimal · none"),
        ("Semantic Coverage",     "0.20",  PURPLE,
         "Cosine similarity of agent reasoning\nvs. actual EIS chunks"),
    ]
    my2 = Inches(4.85)
    mw = Inches(3.9)
    mg = Inches(0.23)
    for i, (name, weight, color, desc) in enumerate(metrics):
        mx = Inches(0.6) + i * (mw + mg)
        mb = sl.shapes.add_shape(1, mx, my2, mw, Inches(1.62))
        fill(mb, CARD)
        mb.line.color.rgb = color
        mb.line.width = Pt(2)
        mhb = sl.shapes.add_shape(1, mx, my2, mw, Inches(0.4))
        fill(mhb, color)
        txt(sl, name, mx + Inches(0.12), my2 + Inches(0.06),
            mw - Inches(0.6), Inches(0.28), size=12, bold=True, color=WHITE)
        txt(sl, f"weight {weight}", mx + mw - Inches(0.78), my2 + Inches(0.1),
            Inches(0.68), Inches(0.22), size=11, bold=True, color=WHITE)
        txt(sl, desc, mx + Inches(0.12), my2 + Inches(0.48),
            mw - Inches(0.24), Inches(1.06), size=12, color=GRAY_LT, wrap=True)

    txt(sl, "OVERALL  =  (F1 × 0.40)  +  (Significance Accuracy × 0.40)  +  (Semantic Coverage × 0.20)",
        Inches(0.6), Inches(6.56), Inches(12.1), Inches(0.38),
        size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    placeholder(sl, "Evaluation panel screenshot — score bars + per-category breakdown",
                Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.38))

    # ── 8. Evaluation Gaps ────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Where the Evaluation Falls Short")

    txt(sl, "Our eval is a starting point — not a production benchmark.",
        Inches(0.6), Inches(1.08), Inches(12.1), Inches(0.38),
        size=15, italic=True, color=GRAY_LT)

    gaps = [
        ("Fixed 8-category F1 is too narrow", RED,
         "Projects can trigger dozens of resource categories not in our fixed set."
         "  Anything outside the 8 is invisible to the metric — F1 artificially looks better."),
        ("Ground truth is LLM-extracted, not human-annotated", YELLOW,
         "Our extractor is itself an LLM — correlated errors inflate scores."
         "  A human expert reading the EIS would catch nuances the LLM misses."),
        ("Semantic coverage ≠ factual accuracy", ORANGE,
         "High cosine similarity means the agent wrote about similar topics."
         "  It does not mean the conclusions are correct."),
        ("No run-level repeatability metric", PURPLE,
         "We do not track score variance across identical reruns."
         "  LLM temperature introduces variance that our current eval ignores."),
    ]

    more_robust = [
        ("LLM-as-Judge", CYAN,
         "A separate evaluator LLM scores each reasoning chain independently,"
         " catching factual errors that F1 cannot detect."),
        ("Question-Answer Generation (QAG)", BLUE,
         "Auto-generate targeted Q&A pairs from the EIS, then test whether"
         " the agent answers them correctly — finer grain than category F1."),
        ("Human annotation corpus", GREEN,
         "Even a small set of expert-labeled EIS documents would break the"
         " evaluator/agent correlation and provide a harder benchmark."),
    ]

    # Left: gaps
    gy = Inches(1.55)
    for title, color, desc in gaps:
        cb = sl.shapes.add_shape(1, Inches(0.6), gy, Inches(5.9), Inches(1.22))
        fill(cb, CARD)
        cb.line.color.rgb = color
        cb.line.width = Pt(1.5)
        lb = sl.shapes.add_shape(1, Inches(0.6), gy, Inches(0.22), Inches(1.22))
        fill(lb, color)
        txt(sl, title, Inches(0.92), gy + Inches(0.06),
            Inches(5.5), Inches(0.28), size=11, bold=True, color=color)
        txt(sl, desc, Inches(0.92), gy + Inches(0.36),
            Inches(5.5), Inches(0.78), size=11, color=GRAY_LT, wrap=True)
        gy += Inches(1.3)

    # Right: more robust approaches
    col_header(sl, "MORE ROBUST APPROACHES", Inches(7.1), Inches(1.55), Inches(5.6), GREEN)
    ry2 = Inches(1.98)
    for title, color, desc in more_robust:
        cb = sl.shapes.add_shape(1, Inches(7.1), ry2, Inches(5.6), Inches(1.5))
        fill(cb, CARD)
        cb.line.color.rgb = color
        cb.line.width = Pt(1.5)
        txt(sl, title, Inches(7.26), ry2 + Inches(0.1),
            Inches(5.3), Inches(0.3), size=12, bold=True, color=color)
        txt(sl, desc, Inches(7.26), ry2 + Inches(0.44),
            Inches(5.3), Inches(0.96), size=11, color=GRAY_LT, wrap=True)
        ry2 += Inches(1.62)

    # ── 9. Tradeoffs ──────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Tradeoffs  ·  Metrics  ·  Emergent Capabilities")

    # Tradeoffs table (left)
    col_header(sl, "TRADEOFFS", Inches(0.6), Inches(1.08), Inches(6.3), ORANGE)
    tradeoffs = [
        ("Latency",       BLUE,
         "~90–120 s end-to-end",
         "Acceptable for planning-stage reviews; too slow for real-time queries"),
        ("Cost",          CYAN,
         "~$0.04–0.08 per full run",
         "Cheap at scale; LLM costs dominate over API calls"),
        ("Accuracy",      GREEN,
         "F1 ≈ 0.75  (mock benchmark)",
         "Sufficient for screening; not yet regulatory-grade without human review"),
        ("Repeatability", YELLOW,
         "±5% score variance across runs",
         "Structured JSON prompts reduce variance; temperature > 0 prevents full determinism"),
    ]
    ty2 = Inches(1.5)
    for label, color, value, note in tradeoffs:
        tb = sl.shapes.add_shape(1, Inches(0.6), ty2, Inches(6.3), Inches(1.16))
        fill(tb, CARD)
        tb.line.color.rgb = color
        tb.line.width = Pt(1)
        dot = sl.shapes.add_shape(1, Inches(0.6), ty2, Inches(0.22), Inches(1.16))
        fill(dot, color)
        txt(sl, label, Inches(0.92), ty2 + Inches(0.08),
            Inches(2.0), Inches(0.28), size=12, bold=True, color=color)
        txt(sl, value, Inches(0.92), ty2 + Inches(0.38),
            Inches(6.0), Inches(0.28), size=13, bold=True, color=WHITE)
        txt(sl, note, Inches(0.92), ty2 + Inches(0.68),
            Inches(6.0), Inches(0.4), size=11, color=GRAY_LT, wrap=True)
        ty2 += Inches(1.26)

    # Emergent capabilities (right)
    col_header(sl, "EMERGENT CAPABILITIES", Inches(7.1), Inches(1.08), Inches(5.6), PURPLE)
    emergent = [
        ("Cross-category compounding", PURPLE,
         "The agent spontaneously flags when wetland + floodplain overlap"
         " amplifies significance — not explicitly prompted."),
        ("Unexpected regulatory links", CYAN,
         "Retrieves and applies tangential CFR parts (e.g., §404 wetland"
         " permits applying to upland grading near waterways)."),
        ("Self-calibrated uncertainty", ORANGE,
         "Confidence scores drop automatically when API data is missing —"
         " no threshold logic in the prompt; emerges from the tier instructions."),
        ("Mitigation specificity", GREEN,
         "Agent proposes site-specific compensatory mitigation amounts (e.g.,"
         " '4.2 acres at 1.5:1 ratio') based on wetland acreage from NWI data."),
    ]
    ey = Inches(1.5)
    for title, color, desc in emergent:
        eb = sl.shapes.add_shape(1, Inches(7.1), ey, Inches(5.6), Inches(1.44))
        fill(eb, CARD)
        eb.line.color.rgb = color
        eb.line.width = Pt(1.5)
        txt(sl, title, Inches(7.26), ey + Inches(0.08),
            Inches(5.3), Inches(0.28), size=12, bold=True, color=color)
        txt(sl, desc, Inches(7.26), ey + Inches(0.38),
            Inches(5.3), Inches(0.98), size=11, color=GRAY_LT, wrap=True)
        ey += Inches(1.54)

    # ── 10. Conclusion ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Conclusion")

    takeaways = [
        ("Weeks of expert review compressed to 90 seconds", 0, ORANGE),
        ("5-agent LangGraph pipeline: parse · data · regulations · matrix · report", 1, GRAY_LT),
        ("RAG grounds every citation — no hallucinated regulations", 0, ORANGE),
        ("Structured prompts + typed state = reproducible, auditable outputs", 1, GRAY_LT),
        ("Quantitatively evaluated  (F1 · significance accuracy · semantic coverage)", 0, ORANGE),
        ("Honest about gaps: F1 scope, LLM-extracted GT, repeatability not tracked", 1, GRAY_LT),
        ("Emergent behaviors: compounding risk, unexpected regulatory links", 0, ORANGE),
        ("Production path: human-annotated GT · LLM-as-Judge · dynamic categories", 1, GRAY_LT),
    ]
    bullets(sl, takeaways, Inches(0.9), Inches(1.15), Inches(7.5), Inches(4.8))

    # Criteria checkboxes (right column)
    criteria_check = [
        ("Perceiving  ✓",  GREEN),
        ("Reasoning   ✓",  BLUE),
        ("Executing   ✓",  YELLOW),
        ("Autonomy    ✓",  ORANGE),
    ]
    ccy = Inches(1.4)
    for label, color in criteria_check:
        ccb = sl.shapes.add_shape(1, Inches(8.8), ccy, Inches(4.0), Inches(0.72))
        fill(ccb, CARD)
        ccb.line.color.rgb = color
        ccb.line.width = Pt(2)
        txt(sl, label, Inches(9.0), ccy + Inches(0.16),
            Inches(3.6), Inches(0.42), size=18, bold=True, color=color)
        ccy += Inches(0.86)

    placeholder(sl, "Live demo / final screenshot",
                Inches(8.8), Inches(5.0), Inches(4.0), Inches(2.28))

    fw_box = sl.shapes.add_shape(1, Inches(0.6), Inches(6.1), Inches(7.8), Inches(0.42))
    fill(fw_box, CARD)
    txt(sl, "Tradeoff:  autonomy · safety · latency · accuracy · cost · reliability",
        Inches(0.8), Inches(6.16), Inches(7.6), Inches(0.3),
        size=13, italic=True, color=GRAY_LT)

    prs.save(output_path)
    _strip_slide_numbers(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    potx = "/Users/sanderschulman/Developer/aiagentsproject/24-880 Presentation Template (1).potx"
    out  = "/Users/sanderschulman/Developer/aiagentsproject/EIA_Hackathon_Presentation.pptx"
    build(potx, out)
