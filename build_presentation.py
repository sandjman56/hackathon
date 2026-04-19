"""Build EIA Agent System presentation following the 24-880 template theme."""
import io
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colors ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x1F, 0x28, 0x33)
ORANGE    = RGBColor(0x00, 0xCC, 0x6A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT   = RGBColor(0xD1, 0xD5, 0xDB)
RED       = RGBColor(0xEF, 0x44, 0x44)
YELLOW    = RGBColor(0xEA, 0xB3, 0x08)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
CYAN      = RGBColor(0x06, 0xB6, 0xD3)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
CARD      = RGBColor(0x2D, 0x37, 0x48)
DARK      = RGBColor(0x1A, 0x20, 0x2C)
DIVIDER   = RGBColor(0x3D, 0x4A, 0x5C)

W = Inches(13.33)
H = Inches(7.5)


# ── Core helpers ─────────────────────────────────────────────────────────────

def convert_potx_to_pptx(potx_path: str) -> Presentation:
    import re as _re
    buf = io.BytesIO()
    with zipfile.ZipFile(potx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename.startswith('ppt/slides/'):
                    continue
                data = zin.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    data = data.replace(
                        b'presentationml.template.main+xml',
                        b'presentationml.presentation.main+xml',
                    )
                    data = _re.sub(
                        rb'<Override[^>]*/ppt/slides/slide\d+\.xml[^>]*/?>',
                        b'', data,
                    )
                if item.filename == 'ppt/presentation.xml':
                    data = _re.sub(rb'<p:sldIdLst>.*?</p:sldIdLst>', b'<p:sldIdLst/>',
                                   data, flags=_re.DOTALL)
                zout.writestr(item, data)
    buf.seek(0)
    prs = Presentation(buf)
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, left, top, width, height, color=CARD, border=None, border_w=1.5):
    s = slide.shapes.add_shape(1, left, top, width, height)
    fill(s, color)
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, size=28):
    """Orange top bar + title text + divider line."""
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
    fill(s, ORANGE)
    txt(slide, title, Inches(0.6), Inches(0.16), Inches(12.5), Inches(0.72),
        size=size, bold=True, color=ORANGE)
    d = slide.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.025))
    fill(d, DIVIDER)


def bullets(slide, items, left, top, width, height):
    """
    items: list of (text, level, color)
    level 0 = main bullet (▸, 18pt bold), level 1 = sub (·, 14pt)
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6 if level == 0 else 2)
        indent = "    " * level
        char   = "▸ " if level == 0 else "· "
        r = p.add_run()
        r.text = indent + char + text
        r.font.size  = Pt(18 if level == 0 else 14)
        r.font.bold  = (level == 0)
        r.font.color.rgb = color


def placeholder(slide, label, left, top, width, height):
    b = slide.shapes.add_shape(1, left, top, width, height)
    fill(b, CARD)
    b.line.color.rgb = ORANGE
    b.line.width = Pt(1.5)
    txt(slide, f"[ {label} ]",
        left + Inches(0.15), top + height / 2 - Pt(24),
        width - Inches(0.3), Inches(0.5),
        size=12, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)


def col_header(slide, label, left, top, width, color=ORANGE):
    b = slide.shapes.add_shape(1, left, top, width, Inches(0.36))
    fill(b, CARD)
    txt(slide, label, left + Inches(0.12), top + Inches(0.05),
        width - Inches(0.24), Inches(0.28),
        size=12, bold=True, color=color)


# ── Build ─────────────────────────────────────────────────────────────────────

def build(potx_path: str, output_path: str):
    prs = convert_potx_to_pptx(potx_path)

    # ── 1. Title ─────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    top_bar = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.8))
    fill(top_bar, ORANGE)
    txt(sl, "EIA AGENT SYSTEM",
        Inches(0.7), Inches(0.2), Inches(11.5), Inches(1.1),
        size=44, bold=True, color=WHITE)
    txt(sl, "Automated Environmental Impact Assessment via Multi-Agent AI",
        Inches(0.7), Inches(2.1), Inches(11.0), Inches(0.7),
        size=22, color=GRAY_LT)
    txt(sl, "24-880  ·  AI Agents for Engineers  ·  Carnegie Mellon University",
        Inches(0.7), Inches(3.05), Inches(10.0), Inches(0.5),
        size=16, color=ORANGE)
    txt(sl, "[ Group Member Names ]",
        Inches(0.7), Inches(3.8), Inches(10.0), Inches(0.4),
        size=15, color=GRAY_LT)
    txt(sl, "[ Presentation Date ]",
        Inches(0.7), Inches(4.3), Inches(10.0), Inches(0.4),
        size=15, color=GRAY_LT)
    bot = sl.shapes.add_shape(1, Inches(0), H - Inches(0.15), W, Inches(0.15))
    fill(bot, ORANGE)

    # ── 2. Background ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Background")

    # Left: NEPA overview
    txt(sl, "What is NEPA?", Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.4),
        size=14, bold=True, color=ORANGE)
    bullets(sl, [
        ("National Environmental Policy Act, signed 1970", 0, WHITE),
        ("Requires agencies to assess environmental impact", 1, GRAY_LT),
        ("before approving major federal projects", 1, GRAY_LT),
        ("Screening-level EA: does this need a full EIS?", 0, WHITE),
        ("Full EIS: years of analysis, millions in cost", 1, GRAY_LT),
    ], Inches(0.6), Inches(1.55), Inches(5.8), Inches(4.5))

    # Right: 5 API source chips
    txt(sl, "Federal Data Sources Queried", Inches(7.0), Inches(1.1),
        Inches(5.9), Inches(0.4), size=14, bold=True, color=ORANGE)
    sources = [
        ("USFWS IPaC",   "Threatened & endangered species",  YELLOW),
        ("Natl. Wetlands Inventory", "Wetland presence & type", CYAN),
        ("FEMA Flood Layer", "100-yr floodplain boundaries", BLUE),
        ("USDA Soil Survey", "Prime & unique farmland",      GREEN),
        ("EPA EJScreen",  "Environmental justice indicators", PURPLE),
    ]
    sy = Inches(1.6)
    for name, desc, color in sources:
        chip = sl.shapes.add_shape(1, Inches(7.0), sy, Inches(5.9), Inches(0.82))
        fill(chip, CARD)
        chip.line.color.rgb = color
        chip.line.width = Pt(1.2)
        txt(sl, name, Inches(7.15), sy + Inches(0.08),
            Inches(5.6), Inches(0.3), size=13, bold=True, color=color)
        txt(sl, desc, Inches(7.15), sy + Inches(0.4),
            Inches(5.6), Inches(0.32), size=12, color=GRAY_LT)
        sy += Inches(0.95)

    placeholder(sl, "NEPA process diagram",
                Inches(0.6), Inches(5.6), Inches(5.8), Inches(1.6))

    # ── 3. Problem Statement ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Problem Statement")

    problems = [
        ("Reviews take weeks of expert analyst time", 0, WHITE),
        ("5+ databases · hundreds of regulations · written report", 1, GRAY_LT),
        ("Existing tools are black boxes", 0, WHITE),
        ("Regulators require traceable citations and reasoning", 1, GRAY_LT),
        ("No objective quality metric for AI-generated outputs", 0, WHITE),
        ("How do we know the agent got it right?", 1, GRAY_LT),
    ]
    bullets(sl, problems, Inches(0.9), Inches(1.15), Inches(11.5), Inches(4.2))

    # Goal statement
    goal = sl.shapes.add_shape(1, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9))
    fill(goal, CARD)
    goal.line.color.rgb = ORANGE
    goal.line.width = Pt(1.5)
    txt(sl, "Goal:  automate EIA screening with explainable, confidence-aware AI —"
           "  fast enough for early-stage planning, rigorous enough for regulatory review.",
        Inches(0.85), Inches(5.7), Inches(11.7), Inches(0.7),
        size=15, color=WHITE)

    # ── 4. Methodology — Architecture ────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — System Architecture")

    # 5 agent boxes
    agents = [
        ("1\nProject\nParser",       BLUE),
        ("2\nEnviron.\nData",        CYAN),
        ("3\nRegulatory\nScreening", GREEN),
        ("4\nImpact\nAnalysis",      ORANGE),
        ("5\nReport\nSynthesis",     YELLOW),
    ]
    bw = Inches(2.22)
    bg_gap = Inches(0.22)
    ax = Inches(0.5)
    ay = Inches(1.18)
    ah = Inches(2.0)
    for i, (name, color) in enumerate(agents):
        bx = ax + i * (bw + bg_gap)
        b = sl.shapes.add_shape(1, bx, ay, bw, ah)
        fill(b, CARD)
        b.line.color.rgb = color
        b.line.width = Pt(2)
        hb = sl.shapes.add_shape(1, bx, ay, bw, Inches(0.42))
        fill(hb, color)
        txt(sl, name, bx + Inches(0.1), ay + Inches(0.48),
            bw - Inches(0.2), ah - Inches(0.58),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    arr = sl.shapes.add_shape(1, Inches(0.6), ay + ah + Inches(0.12),
                              Inches(12.1), Inches(0.03))
    fill(arr, DIVIDER)
    txt(sl, "sequential  ▶", Inches(5.5), ay + ah + Inches(0.18),
        Inches(3.0), Inches(0.3), size=10, italic=True, color=DIVIDER)

    # Stack row
    stack = [
        ("React + Vite",  "SSE streaming UI — real-time status per agent",    BLUE),
        ("FastAPI",       "POST /api/run  →  SSE generator  →  PostgreSQL",   CYAN),
        ("LangGraph",     "Typed state flows through sequential graph nodes",  GREEN),
        ("pgvector",      "HNSW cosine index for RAG + evaluation search",    ORANGE),
    ]
    sy = Inches(3.55)
    sw = Inches(12.1)
    sh = Inches(0.6)
    for i, (label, desc, color) in enumerate(stack):
        rb = sl.shapes.add_shape(1, Inches(0.6), sy + i * (sh + Inches(0.08)),
                                 sw, sh)
        fill(rb, CARD)
        rb.line.color.rgb = color
        rb.line.width = Pt(1)
        txt(sl, label,
            Inches(0.8), sy + i * (sh + Inches(0.08)) + Inches(0.12),
            Inches(1.9), Inches(0.36), size=13, bold=True, color=color)
        txt(sl, desc,
            Inches(2.9), sy + i * (sh + Inches(0.08)) + Inches(0.14),
            Inches(9.7), Inches(0.32), size=12, color=GRAY_LT)

    placeholder(sl, "Pipeline / Brain Scanner screenshot",
                Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.25))

    # ── 5. Agent 1: Project Parser ────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Agent 1: Project Parser")
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.04), Inches(12.1), Inches(0.07))
    fill(d, BLUE)

    txt(sl, "Converts a plain-English project description into structured metadata"
           " used by every downstream agent.  Model: Gemini 2.5 Flash",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.6),
        size=14, color=GRAY_LT)

    # Input card (left)
    col_header(sl, "INPUTS", Inches(0.6), Inches(2.0), Inches(5.8), CYAN)
    bullets(sl, [
        ("Project name", 0, WHITE),
        ("GPS coordinates  (lat, lon)", 0, WHITE),
        ("Free-text project description", 0, WHITE),
    ], Inches(0.6), Inches(2.42), Inches(5.8), Inches(2.5))

    # Output card (right)
    col_header(sl, "OUTPUTS  (stored as JSONB)", Inches(7.0), Inches(2.0), Inches(5.9), BLUE)
    bullets(sl, [
        ("project_type  (highway, solar farm, pipeline…)", 0, WHITE),
        ("scale  (small / medium / large)", 0, WHITE),
        ("location  (city, county, state)", 0, WHITE),
        ("actions  →  column headers in impact matrix", 0, WHITE),
        ("federal_nexus  (triggers NEPA applicability)", 0, WHITE),
    ], Inches(7.0), Inches(2.42), Inches(5.9), Inches(3.0))

    placeholder(sl, "Project Parser output — key/value view in IMPORT RUN panel",
                Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.65))

    # ── 6. Agent 2: Environmental Data ───────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Agent 2: Environmental Data Agent")
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.04), Inches(12.1), Inches(0.07))
    fill(d, CYAN)

    txt(sl, "Queries 5 federal REST APIs using the project's GPS coordinates."
           "  No LLM — pure deterministic API calls.",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
        size=14, color=GRAY_LT)

    apis = [
        ("USFWS IPaC",                 "Threatened & endangered species in project area",    YELLOW),
        ("National Wetlands Inventory","Wetland polygons, type codes, acreage",              CYAN),
        ("FEMA Flood Hazard Layer",    "Zone designations (AE, X…) and SFHA flag",          BLUE),
        ("USDA Web Soil Survey",       "Prime / unique farmland classification",             GREEN),
        ("EPA EJScreen",               "Minority %, low-income %, pollution percentiles",   PURPLE),
    ]
    for i, (name, desc, color) in enumerate(apis):
        aw = Inches(2.26)
        ax2 = Inches(0.5) + i * (aw + Inches(0.2))
        ab = sl.shapes.add_shape(1, ax2, Inches(1.88), aw, Inches(2.5))
        fill(ab, CARD)
        ab.line.color.rgb = color
        ab.line.width = Pt(1.5)
        hb = sl.shapes.add_shape(1, ax2, Inches(1.88), aw, Inches(0.4))
        fill(hb, color)
        txt(sl, f"{i+1}", ax2 + Inches(0.1), Inches(1.92),
            Inches(0.3), Inches(0.32), size=13, bold=True, color=WHITE)
        txt(sl, name, ax2 + Inches(0.1), Inches(2.35),
            aw - Inches(0.2), Inches(0.55), size=12, bold=True, color=color)
        txt(sl, desc, ax2 + Inches(0.1), Inches(2.96),
            aw - Inches(0.2), Inches(1.3), size=11, color=GRAY_LT, wrap=True)

    txt(sl, "Partial failures handled gracefully — missing data lowers confidence score (not a crash).",
        Inches(0.6), Inches(4.52), Inches(12.1), Inches(0.4),
        size=13, italic=True, color=ORANGE)

    placeholder(sl, "Environmental Data output — API response cards in IMPORT RUN panel",
                Inches(0.6), Inches(5.05), Inches(12.1), Inches(2.2))

    # ── 7. Agent 3: Regulatory Screening ─────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Agent 3: Regulatory Screening  (RAG)")
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.04), Inches(12.1), Inches(0.07))
    fill(d, GREEN)

    txt(sl, "Retrieves applicable regulations using vector similarity search"
           " over an indexed corpus of federal CFR parts and state codes.",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
        size=14, color=GRAY_LT)

    # Left: Ingestion
    col_header(sl, "INGESTION  (offline)", Inches(0.6), Inches(1.82), Inches(5.8), CYAN)
    bullets(sl, [
        ("Upload PDFs via /api/regulations/sources", 0, WHITE),
        ("Or fetch live eCFR XML by title + part", 0, WHITE),
        ("Chunk → embed → pgvector (HNSW cosine index)", 0, WHITE),
        ("Sources scoped to a specific project", 0, WHITE),
    ], Inches(0.6), Inches(2.24), Inches(5.8), Inches(2.8))

    # Right: Retrieval
    col_header(sl, "RETRIEVAL  (at run time)", Inches(7.0), Inches(1.82), Inches(5.9), GREEN)
    bullets(sl, [
        ("Embed project context as query vector", 0, WHITE),
        ("Top-K similarity search over project's sources", 0, WHITE),
        ("LLM identifies triggered regulations from chunks", 0, WHITE),
        ("Output: name · citation · jurisdiction · why", 0, WHITE),
    ], Inches(7.0), Inches(2.24), Inches(5.9), Inches(2.8))

    # Model badge
    badge = sl.shapes.add_shape(1, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.42))
    fill(badge, CARD)
    txt(sl, "Model: claude-haiku-4-5   ·   Embeddings: text-embedding-3-small   ·   Vector DB: pgvector",
        Inches(0.8), Inches(4.97), Inches(11.8), Inches(0.3),
        size=12, color=GRAY_LT)

    placeholder(sl, "Regulatory Screening output — regulation cards in IMPORT RUN panel",
                Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.8))

    # ── 8. Agent 4: Impact Analysis ───────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Agent 4: Impact Analysis")
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.04), Inches(12.1), Inches(0.07))
    fill(d, ORANGE)

    txt(sl, "Populates a significance matrix: every project action × every environmental"
           " resource category, governed by the applicable regulation.",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.55),
        size=14, color=GRAY_LT)

    # Matrix sample (left)
    rows = [
        ("Category",         "Site Prep",   "Construction", "Operation"),
        ("Wetlands",         "moderate",    "significant",  "minimal"),
        ("Endangered Spp.",  "none",        "significant",  "none"),
        ("Floodplain",       "moderate",    "moderate",     "none"),
        ("Air Quality",      "minimal",     "moderate",     "minimal"),
    ]
    sig_colors = {
        "Category": ORANGE, "Site Prep": ORANGE, "Construction": ORANGE, "Operation": ORANGE,
        "significant": RED, "moderate": YELLOW, "minimal": GREEN, "none": CYAN,
    }
    cw = Inches(1.82)
    ch = Inches(0.52)
    mx = Inches(0.6)
    my = Inches(1.92)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cx = mx + c * cw
            cy = my + r * ch
            cb = sl.shapes.add_shape(1, cx, cy, cw - Inches(0.03), ch - Inches(0.025))
            is_hdr = r == 0 or c == 0
            fill(cb, DARK if is_hdr else CARD)
            cb.line.color.rgb = DIVIDER
            cb.line.width = Pt(0.75)
            txt(sl, val, cx + Inches(0.07), cy + Inches(0.08),
                cw - Inches(0.16), ch - Inches(0.14),
                size=11, bold=is_hdr, color=sig_colors.get(val, WHITE))

    # Legend
    legend = [("significant", RED), ("moderate", YELLOW), ("minimal", GREEN), ("none", CYAN)]
    lx = Inches(0.6)
    ly = my + len(rows) * ch + Inches(0.12)
    txt(sl, "Scale:", lx, ly, Inches(0.8), Inches(0.28), size=11, bold=True, color=WHITE)
    for i, (label, color) in enumerate(legend):
        dot = sl.shapes.add_shape(1, lx + Inches(0.9) + i * Inches(1.7),
                                  ly + Inches(0.05), Inches(0.16), Inches(0.16))
        fill(dot, color)
        txt(sl, label, lx + Inches(1.12) + i * Inches(1.7), ly,
            Inches(1.5), Inches(0.28), size=11, color=color)

    # Per-cell fields (right)
    txt(sl, "Each cell contains:", Inches(8.0), Inches(1.88), Inches(5.0), Inches(0.32),
        size=13, bold=True, color=ORANGE)
    fields = [
        ("significance", "significant / moderate / minimal / none", ORANGE),
        ("confidence",   "0.0 – 1.0  (see next slide)",            YELLOW),
        ("reasoning",    "1–2 sentence explanation",                GRAY_LT),
        ("mitigation",   "avoidance / minimization / compensatory", GREEN),
        ("needs_review", "True when confidence < 0.60",             RED),
    ]
    fy = Inches(2.28)
    for fname, fdesc, fcolor in fields:
        fb = sl.shapes.add_shape(1, Inches(8.0), fy, Inches(5.1), Inches(0.68))
        fill(fb, CARD)
        fb.line.color.rgb = fcolor
        fb.line.width = Pt(0.8)
        txt(sl, fname, Inches(8.15), fy + Inches(0.05),
            Inches(4.8), Inches(0.28), size=12, bold=True, color=fcolor)
        txt(sl, fdesc, Inches(8.15), fy + Inches(0.33),
            Inches(4.8), Inches(0.28), size=11, color=GRAY_LT)
        fy += Inches(0.76)

    placeholder(sl, "Impact matrix — color-coded significance table from a real run",
                Inches(0.6), Inches(5.95), Inches(7.0), Inches(1.3))

    # ── 9. Confidence Score Tiers ─────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Confidence Score Tiers  (Impact Analysis Agent)")

    txt(sl, "The LLM assigns a confidence score to each matrix cell based on"
           " data availability and regulatory specificity.",
        Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
        size=14, color=GRAY_LT)

    # Column headers
    col_xs = [Inches(0.6), Inches(2.3), Inches(4.2), Inches(6.0)]
    col_ws = [Inches(1.6), Inches(1.8), Inches(1.7), Inches(6.9)]
    col_labels = ["SCORE", "TIER", "DATA?", "WHAT IT MEANS"]
    hdr_row = sl.shapes.add_shape(1, Inches(0.6), Inches(1.72), Inches(12.1), Inches(0.38))
    fill(hdr_row, DARK)
    for cx, cw2, cl in zip(col_xs, col_ws, col_labels):
        txt(sl, cl, cx + Inches(0.1), Inches(1.76),
            cw2 - Inches(0.1), Inches(0.28), size=11, bold=True, color=ORANGE)

    # Tier rows — from impact_analysis.py system prompt
    tiers = [
        ("0.85–1.0",  "HIGH",     GREEN,   "Yes — quantified", "API returned specific data AND regulation has explicit numeric thresholds"),
        ("0.65–0.84", "GOOD",     CYAN,    "Yes — general",    "API data present but regulatory threshold is a judgment call (no explicit number)"),
        ("0.45–0.64", "MODERATE", YELLOW,  "Partial",          "Some APIs errored / empty, or regulation is tangentially related"),
        ("0.25–0.44", "LOW",      ORANGE,  "No",               "No direct data — determination relies on general domain knowledge"),
        ("0.00–0.24", "VERY LOW", RED,     "No",               "Pure extrapolation — no supporting data or regulatory context"),
    ]
    row_h = Inches(0.72)
    for i, (score, tier, color, data, meaning) in enumerate(tiers):
        ry = Inches(2.1) + i * row_h
        rb = sl.shapes.add_shape(1, Inches(0.6), ry, Inches(12.1), row_h - Inches(0.04))
        fill(rb, CARD)
        rb.line.color.rgb = DIVIDER
        rb.line.width = Pt(0.5)
        # Score range
        txt(sl, score, Inches(0.72), ry + Inches(0.16),
            Inches(1.5), Inches(0.4), size=15, bold=True, color=color)
        # Tier
        txt(sl, tier, Inches(2.35), ry + Inches(0.18),
            Inches(1.6), Inches(0.36), size=13, bold=True, color=color)
        # Data?
        txt(sl, data, Inches(4.22), ry + Inches(0.18),
            Inches(1.7), Inches(0.36), size=12, color=GRAY_LT)
        # Meaning
        txt(sl, meaning, Inches(6.05), ry + Inches(0.1),
            Inches(6.5), Inches(0.54), size=12, color=GRAY_LT, wrap=True)

    # needs_review flag
    flag = sl.shapes.add_shape(1, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.48))
    fill(flag, RGBColor(0x3B, 0x1A, 0x1A))
    flag.line.color.rgb = RED
    flag.line.width = Pt(1.2)
    txt(sl, "needs_review = True  when  confidence < 0.60"
           "  →  cell is highlighted in the UI for human expert review",
        Inches(0.8), Inches(5.86), Inches(11.7), Inches(0.36),
        size=14, bold=True, color=RED)

    placeholder(sl, "Impact matrix UI — needs_review cells highlighted",
                Inches(0.6), Inches(6.42), Inches(12.1), Inches(0.85))

    # ── 10. Agent 5: Report Synthesis ─────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Agent 5: Report Synthesis")
    d = sl.shapes.add_shape(1, Inches(0.6), Inches(1.04), Inches(12.1), Inches(0.07))
    fill(d, YELLOW)

    txt(sl, "Generates a 10-section NEPA Environmental Assessment narrative."
           "  Sections 1–7: LLM-written with inline citations.  Sections 8–10: templated.",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.55),
        size=14, color=GRAY_LT)

    # Sections list — two columns
    left_secs = [
        "1.  Project Description",
        "2.  Purpose and Need",
        "3.  Alternatives Analysis",
        "4.  Affected Environment",
        "5.  Environmental Consequences",
    ]
    right_secs = [
        "6.  Cumulative Impacts",
        "7.  Mitigation Measures",
        "8.  Regulatory Compliance Summary",
        "9.  Agency Coordination",
        "10. References & Data Sources",
    ]
    for secs, lx2 in [(left_secs, Inches(0.6)), (right_secs, Inches(7.0))]:
        tb = sl.shapes.add_textbox(lx2, Inches(1.9), Inches(5.8), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for s in secs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(6)
            r = p.add_run()
            r.text = s
            r.font.size = Pt(17)
            r.font.color.rgb = WHITE

    badge2 = sl.shapes.add_shape(1, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.42))
    fill(badge2, CARD)
    txt(sl, "Model: Gemini 2.5 Flash   ·   Every statement cites the regulation and federal data source that drove it.",
        Inches(0.8), Inches(4.77), Inches(11.8), Inches(0.3), size=12, color=GRAY_LT)

    placeholder(sl, "Report Synthesis output — sectioned narrative in IMPORT RUN panel",
                Inches(0.6), Inches(5.28), Inches(12.1), Inches(2.0))

    # ── 11. Methodology — Evaluation Process ──────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Methodology — Evaluation Process")

    txt(sl, "Upload a published EIS document for the same project."
           "  Score the agent's impact matrix against it.",
        Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
        size=14, color=GRAY_LT)

    steps = [
        ("UPLOAD",   "User uploads a\npublished EIS PDF\nfor a saved project",   BLUE),
        ("PARSE",    "PDF chunked,\nembedded, stored\nin pgvector",               CYAN),
        ("EXTRACT",  "LLM reads chunks\n→ ground truth:\ncategory + significance", GREEN),
        ("SCORE",    "Agent matrix vs.\nground truth →\n3 metrics computed",      ORANGE),
        ("REVIEW",   "Per-category\nTP / FP / FN\nbreakdown shown",              YELLOW),
    ]
    sw2 = Inches(2.3)
    sh2 = Inches(3.6)
    sg  = Inches(0.22)
    sx2 = Inches(0.5)
    sy2 = Inches(1.78)
    for i, (label, desc, color) in enumerate(steps):
        bx2 = sx2 + i * (sw2 + sg)
        sb = sl.shapes.add_shape(1, bx2, sy2, sw2, sh2)
        fill(sb, CARD)
        sb.line.color.rgb = color
        sb.line.width = Pt(1.5)
        hb2 = sl.shapes.add_shape(1, bx2, sy2, sw2, Inches(0.42))
        fill(hb2, color)
        txt(sl, f"{i+1}. {label}", bx2 + Inches(0.1), sy2 + Inches(0.07),
            sw2 - Inches(0.2), Inches(0.32), size=11, bold=True, color=WHITE)
        txt(sl, desc, bx2 + Inches(0.12), sy2 + Inches(0.52),
            sw2 - Inches(0.24), sh2 - Inches(0.62),
            size=13, color=GRAY_LT, wrap=True)

    txt(sl, "Ground truth extraction is cached — runs once per EIS document, not per scoring request.",
        Inches(0.6), Inches(5.48), Inches(12.1), Inches(0.38),
        size=13, italic=True, color=ORANGE)

    placeholder(sl, "Evaluation page — upload panel + score bars",
                Inches(0.6), Inches(5.98), Inches(12.1), Inches(1.3))

    # ── 12. Prompts (1/2) — Agents 1–3 + Evaluation ──────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Prompts  (1/2)  —  15 total across 5 LLM callers", size=24)

    txt(sl, "Every agent uses a system prompt (role + rules) paired with a user prompt"
           " (runtime data). All non-report prompts enforce JSON-only output.",
        Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.48),
        size=13, color=GRAY_LT)

    # Table header
    hrow = sl.shapes.add_shape(1, Inches(0.6), Inches(1.68), Inches(12.1), Inches(0.34))
    fill(hrow, DARK)
    for label, lx2 in [("CALLER", Inches(0.72)), ("SYSTEM PROMPT", Inches(2.35)),
                       ("USER PROMPT", Inches(6.15)), ("FILE", Inches(9.95))]:
        txt(sl, label, lx2, Inches(1.72), Inches(3.7), Inches(0.26),
            size=10, bold=True, color=ORANGE)

    prompt_rows = [
        (
            "Project Parser", BLUE,
            "Extract structured metadata from NL input; respond JSON only",
            "Format name + coords + description → 5 JSON fields\n(project_type, scale, location, actions, federal_nexus)",
            "agents/project_parser.py",
        ),
        (
            "Regulatory Screening", GREEN,
            "NEPA compliance assistant; identify permits & approvals;\nno invented citations; JSON array only",
            "RAG: top-8 retrieved CFR chunks + project context\n→ list of triggered regulations with citations",
            "agents/regulatory_screening.py",
        ),
        (
            "Impact Analysis", ORANGE,
            "Assess every action × category × regulation;\nconfidence tiers tied to data quality (see slide 9)",
            "Compile env API data + regulations + actions\n→ cells JSON with significance + confidence + mitigation",
            "agents/impact_analysis.py",
        ),
        (
            "GT Extractor\n(Evaluation)", PURPLE,
            "Expert EIS analyst; extract ALL resource categories\nand significance levels; JSON array only",
            "Up to 40 sampled EIS chunks (tables prioritized)\n→ ground truth: category + significance + evidence",
            "rag_eval/extractor.py",
        ),
    ]

    row_h = Inches(1.12)
    for i, (name, color, sys_txt, usr_txt, fpath) in enumerate(prompt_rows):
        ry = Inches(2.08) + i * row_h
        rb = sl.shapes.add_shape(1, Inches(0.6), ry, Inches(12.1), row_h - Inches(0.06))
        fill(rb, CARD if i % 2 == 0 else BG)
        rb.line.color.rgb = DIVIDER
        rb.line.width = Pt(0.5)

        # Caller name chip
        chip = sl.shapes.add_shape(1, Inches(0.62), ry + Inches(0.08),
                                   Inches(1.6), row_h - Inches(0.22))
        fill(chip, DARK)
        chip.line.color.rgb = color
        chip.line.width = Pt(1.2)
        txt(sl, name, Inches(0.72), ry + Inches(0.18),
            Inches(1.42), Inches(0.72), size=11, bold=True, color=color)

        # System prompt
        txt(sl, sys_txt, Inches(2.35), ry + Inches(0.1),
            Inches(3.72), row_h - Inches(0.2), size=11, color=WHITE, wrap=True)

        # User prompt
        txt(sl, usr_txt, Inches(6.15), ry + Inches(0.1),
            Inches(3.72), row_h - Inches(0.2), size=11, color=GRAY_LT, wrap=True)

        # File path
        txt(sl, fpath, Inches(9.95), ry + Inches(0.28),
            Inches(2.65), Inches(0.5), size=9, italic=True, color=ORANGE)

    # ── 13. Prompts (2/2) — Report Synthesis (7 prompts) ─────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Prompts  (2/2)  —  Agent 5: Report Synthesis  (7 prompts)", size=24)

    txt(sl, "One system prompt sets the NEPA technical-writer role."
           "  Six separate LLM calls generate individual EA sections.",
        Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.42),
        size=13, color=GRAY_LT)

    # System prompt card
    sys_box = sl.shapes.add_shape(1, Inches(0.6), Inches(1.62), Inches(12.1), Inches(0.62))
    fill(sys_box, DARK)
    sys_box.line.color.rgb = YELLOW
    sys_box.line.width = Pt(1.2)
    txt(sl, "SYSTEM  (shared across all 6 section calls)",
        Inches(0.75), Inches(1.66), Inches(5.0), Inches(0.28),
        size=11, bold=True, color=YELLOW)
    txt(sl, "NEPA technical writer · professional passive-voice prose · no invented data"
           " · return section content only (no title/number) · 3–5 sentence paragraphs",
        Inches(0.75), Inches(1.9), Inches(9.5), Inches(0.28),
        size=12, color=GRAY_LT)
    txt(sl, "agents/report_synthesis.py", Inches(10.8), Inches(1.74),
        Inches(1.85), Inches(0.28), size=9, italic=True, color=ORANGE)

    # 6 section prompt cards (2 columns × 3 rows)
    section_prompts = [
        ("§1  Purpose & Need",
         "1–2 paragraphs: why this project,\nwhat need it addresses, where",
         "agents/templates/nepa_ea.py"),
        ("§2  Proposed Action",
         "2–3 paragraphs: describe project\nand each discrete action",
         "agents/templates/nepa_ea.py"),
        ("§3  No-Action Alternative",
         "1 paragraph: what happens if\nproject does not proceed",
         "agents/templates/nepa_ea.py"),
        ("§4  Affected Environment",
         "Convert raw API data into readable\nnarrative of existing conditions",
         "agents/templates/nepa_ea.py"),
        ("§5  Environmental Consequences",
         "Impact per category: actions causing it,\nsignificance level, regulatory basis",
         "agents/templates/nepa_ea.py"),
        ("§6  Mitigation Measures",
         "Describe mitigations grouped by type:\navoidance / minimization / compensatory",
         "agents/templates/nepa_ea.py"),
    ]

    col_w3 = Inches(5.9)
    card_h = Inches(1.3)
    card_g = Inches(0.14)
    for i, (sec_title, sec_desc, fpath2) in enumerate(section_prompts):
        col = i % 2
        row = i // 2
        cx3 = Inches(0.6) + col * (col_w3 + Inches(0.3))
        cy3 = Inches(2.38) + row * (card_h + card_g)
        cb2 = sl.shapes.add_shape(1, cx3, cy3, col_w3, card_h)
        fill(cb2, CARD)
        cb2.line.color.rgb = YELLOW
        cb2.line.width = Pt(0.8)
        txt(sl, sec_title, cx3 + Inches(0.15), cy3 + Inches(0.08),
            col_w3 - Inches(0.3), Inches(0.3), size=12, bold=True, color=YELLOW)
        txt(sl, sec_desc, cx3 + Inches(0.15), cy3 + Inches(0.42),
            col_w3 - Inches(0.3), Inches(0.72), size=12, color=GRAY_LT, wrap=True)
        txt(sl, fpath2, cx3 + col_w3 - Inches(2.65), cy3 + card_h - Inches(0.26),
            Inches(2.55), Inches(0.22), size=8, italic=True, color=ORANGE,
            align=PP_ALIGN.RIGHT)

    # ── 14. Results — Agent Outputs ───────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Results — Agent Outputs")

    out_items = [
        ("Project\nParser",       BLUE,   'project_type: "highway"\nactions: ["grading","paving"]\nfederal_nexus: true'),
        ("Environ.\nData",        CYAN,   "Species: Indiana bat\nWetlands: 4.2 ac PEM1C\nFEMA: Zone AE\nFarmland: Prime"),
        ("Regulatory\nScreening", GREEN,  "CWA §404 — Wetlands\nESA §7 — Species\nEO 11990 — Floodplain"),
        ("Impact\nMatrix",        ORANGE, "wetlands × grading:\n  sig: significant\n  conf: 0.82\n  review: false"),
        ("Report\nSynthesis",     YELLOW, '"4.2 ac of PEM1C wetland\nintersects the 100-yr\nfloodplain (Zone AE)…"'),
    ]
    bw2 = Inches(2.3)
    bh2 = Inches(3.2)
    bg2 = Inches(0.22)
    bx3 = Inches(0.5)
    by3 = Inches(1.18)
    for name, color, sample in out_items:
        ob = sl.shapes.add_shape(1, bx3, by3, bw2, bh2)
        fill(ob, CARD)
        ob.line.color.rgb = color
        ob.line.width = Pt(1.5)
        ohb = sl.shapes.add_shape(1, bx3, by3, bw2, Inches(0.4))
        fill(ohb, color)
        txt(sl, name, bx3 + Inches(0.08), by3 + Inches(0.46),
            bw2 - Inches(0.16), Inches(0.6), size=11, bold=True, color=color)
        txt(sl, sample, bx3 + Inches(0.12), by3 + Inches(1.12),
            bw2 - Inches(0.24), bh2 - Inches(1.22),
            size=10, color=GRAY_LT, wrap=True)
        bx3 += bw2 + bg2

    placeholder(sl, "IMPORT RUN panel — all 5 agents expanded",
                Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.72))

    # ── 13. Results — Evaluation Scores ───────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Results — Evaluation Scores  (mock data)")

    # Score bars (left)
    score_items = [
        ("Overall",                0.78, "78%",  ORANGE),
        ("Category F1",            0.75, "0.75",  BLUE),
        ("Precision",              0.83, "0.83",  CYAN),
        ("Recall",                 0.67, "0.67",  GREEN),
        ("Significance Accuracy",  0.81, "0.81",  YELLOW),
        ("Semantic Coverage",      0.72, "0.72",  PURPLE),
    ]
    bar_x   = Inches(0.6)
    bar_y   = Inches(1.15)
    bar_max = Inches(4.8)
    bar_h2  = Inches(0.52)
    bar_g   = Inches(0.1)
    for label, val, val_str, color in score_items:
        by2 = bar_y + score_items.index((label, val, val_str, color)) * (bar_h2 + bar_g)
        txt(sl, label, bar_x, by2 + Inches(0.1), Inches(2.3), Inches(0.35), size=13, color=WHITE)
        track = sl.shapes.add_shape(1, bar_x + Inches(2.4), by2 + Inches(0.12),
                                    bar_max, bar_h2 - Inches(0.22))
        fill(track, DARK)
        fb = sl.shapes.add_shape(1, bar_x + Inches(2.4), by2 + Inches(0.12),
                                 max(Inches(0.08), bar_max * val), bar_h2 - Inches(0.22))
        fill(fb, color)
        txt(sl, val_str, bar_x + Inches(7.35), by2 + Inches(0.08),
            Inches(0.8), Inches(0.35), size=13, bold=True, color=color)

    # Per-category table
    txt(sl, "PER-CATEGORY BREAKDOWN", Inches(0.6), Inches(5.42),
        Inches(7.8), Inches(0.32), size=12, bold=True, color=ORANGE)
    tbl_headers = ["Category", "Label", "Agent", "Ground Truth"]
    tbl_widths  = [Inches(2.5), Inches(0.85), Inches(1.35), Inches(1.4)]
    tbl_x_starts = [Inches(0.6), Inches(3.15), Inches(4.05), Inches(5.45)]
    tbl_y = Inches(5.8)
    hrow = sl.shapes.add_shape(1, Inches(0.6), tbl_y, Inches(6.95), Inches(0.28))
    fill(hrow, DARK)
    for hdr_txt, hx, hw in zip(tbl_headers, tbl_x_starts, tbl_widths):
        txt(sl, hdr_txt, hx + Inches(0.06), tbl_y + Inches(0.04),
            hw, Inches(0.22), size=10, bold=True, color=ORANGE)

    tbl_rows = [
        ("wetlands",            "TP", "significant", "significant"),
        ("floodplain",          "TP", "moderate",    "moderate"),
        ("endangered_species",  "FP", "significant", "none"),
        ("prime_farmland",      "FN", "none",        "moderate"),
        ("environmental_justice","TP","moderate",    "moderate"),
    ]
    lc = {"TP": GREEN, "FP": RED, "FN": YELLOW, "TN": CYAN}
    for ri, row_data in enumerate(tbl_rows):
        ry2 = tbl_y + Inches(0.28) + ri * Inches(0.2)
        rb2 = sl.shapes.add_shape(1, Inches(0.6), ry2, Inches(6.95), Inches(0.19))
        fill(rb2, CARD if ri % 2 == 0 else BG)
        for j, (cell_v, hx, hw) in enumerate(zip(row_data, tbl_x_starts, tbl_widths)):
            c2 = lc.get(cell_v, GRAY_LT) if j == 1 else (WHITE if j == 0 else GRAY_LT)
            txt(sl, cell_v, hx + Inches(0.06), ry2 + Inches(0.02),
                hw, Inches(0.16), size=9, bold=(j == 1), color=c2)

    placeholder(sl, "Evaluation panel screenshot — score bars + methodology modal",
                Inches(8.1), Inches(1.15), Inches(5.0), Inches(6.1))

    # ── 14. Evaluation Metrics Explained ─────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Results — How Scores Are Calculated", size=26)

    col_w2 = Inches(3.9)
    col_g  = Inches(0.22)
    cx0    = Inches(0.6)
    cy0    = Inches(1.1)

    metric_data = [
        ("CATEGORY F1", "× 0.40", BLUE,
         [
             "8 fixed categories evaluated:",
             "wetlands · air_quality · noise · traffic",
             "env. justice · species · floodplain · farmland",
             "",
             "TP  agent flagged  AND  EIS confirms",
             "FP  agent flagged  BUT  EIS says none",
             "FN  EIS confirms   BUT  agent missed",
             "",
             "Precision = TP / (TP+FP)",
             "Recall    = TP / (TP+FN)",
             "F1 = 2·P·R / (P+R)",
         ]),
        ("SIGNIFICANCE ACCURACY", "× 0.40", GREEN,
         [
             "Ordinal significance scale:",
             "  significant = 3",
             "  moderate    = 2",
             "  minimal     = 1",
             "  none        = 0",
             "",
             "For each matched category:",
             "  |diff| = 0  →  1.0  (exact match)",
             "  |diff| = 1  →  0.5  (off by one)",
             "  |diff| ≥ 2  →  0.0",
             "",
             "Score = mean across all pairs",
         ]),
        ("SEMANTIC COVERAGE", "× 0.20", YELLOW,
         [
             "Measures: does the agent's reasoning",
             "align with the actual EIS text?",
             "",
             "Up to 10 reasoning snippets embedded",
             "→ cosine similarity vs. EIS chunks",
             "",
             "score = mean(max cosine similarity)",
             "        over all reasoning snippets",
             "",
             "No LLM — pure vector math",
             "on pre-stored embeddings",
         ]),
    ]

    for i, (title, weight, color, lines) in enumerate(metric_data):
        bx4 = cx0 + i * (col_w2 + col_g)
        mb = sl.shapes.add_shape(1, bx4, cy0, col_w2, Inches(5.5))
        fill(mb, CARD)
        mb.line.color.rgb = color
        mb.line.width = Pt(2)
        mhb = sl.shapes.add_shape(1, bx4, cy0, col_w2, Inches(0.55))
        fill(mhb, color)
        txt(sl, title, bx4 + Inches(0.1), cy0 + Inches(0.04),
            col_w2 - Inches(0.5), Inches(0.3), size=12, bold=True, color=WHITE)
        txt(sl, weight, bx4 + col_w2 - Inches(0.7), cy0 + Inches(0.18),
            Inches(0.6), Inches(0.28), size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

        tb2 = sl.shapes.add_textbox(bx4 + Inches(0.12), cy0 + Inches(0.62),
                                    col_w2 - Inches(0.24), Inches(4.65))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        first2 = True
        for line in lines:
            p2 = tf2.paragraphs[0] if first2 else tf2.add_paragraph()
            first2 = False
            p2.space_before = Pt(3)
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(12)
            r2.font.color.rgb = GRAY_LT if line else WHITE
            r2.font.bold = any(kw in line for kw in ("Precision", "Recall", "F1", "score ="))

    txt(sl, "OVERALL  =  (F1 × 0.40)  +  (Significance Accuracy × 0.40)  +  (Semantic Coverage × 0.20)",
        Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.45),
        size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ── 15. Discussion ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Discussion")

    col_header(sl, "CURRENT SHORTCOMINGS", Inches(0.6), Inches(1.1), Inches(5.8), RED)
    bullets(sl, [
        ("Ground truth is LLM-extracted, not annotated by humans", 0, WHITE),
        ("LLM misread of EIS → inflated scores", 1, GRAY_LT),
        ("Only 8 fixed categories scored (F1 scope)", 0, WHITE),
        ("Project-specific impacts invisible to metric", 1, GRAY_LT),
        ("Semantic similarity ≠ factual accuracy", 0, WHITE),
        ("Style match can score high without correct reasoning", 1, GRAY_LT),
        ("LLM non-determinism — same run, different score", 0, WHITE),
    ], Inches(0.6), Inches(1.52), Inches(5.8), Inches(5.4))

    col_header(sl, "FUTURE IMPROVEMENTS", Inches(7.0), Inches(1.1), Inches(5.9), GREEN)
    bullets(sl, [
        ("Human-annotated ground truth corpus", 0, WHITE),
        ("Domain experts review LLM extractions", 1, GRAY_LT),
        ("Dynamic category set per project type", 0, WHITE),
        ("Replace 8 fixed with project-driven taxonomy", 1, GRAY_LT),
        ("LLM-as-Judge for reasoning quality", 0, WHITE),
        ("Separate evaluator scores each reasoning chain", 1, GRAY_LT),
        ("Run-level repeatability tracking", 0, WHITE),
    ], Inches(7.0), Inches(1.52), Inches(5.9), Inches(5.4))

    # ── 16. Conclusion ────────────────────────────────────────────────────────
    # (keeping as slide 15 since we now have exactly 15 slots used)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, "Conclusion")

    takeaways = [
        ("Weeks → minutes  for NEPA screening-level review", 0, ORANGE),
        ("5 agents: parse · data · regulations · matrix · report", 1, GRAY_LT),
        ("Every cell cites the regulation and data behind it", 1, GRAY_LT),
        ("Confidence-aware — flags uncertainty for human review", 0, ORANGE),
        ("5-tier model tied to actual data quality (not vibes)", 1, GRAY_LT),
        ("Quantitatively evaluated against real EIS documents", 0, ORANGE),
        ("F1 + Significance Accuracy + Semantic Coverage", 1, GRAY_LT),
        ("Per-category TP/FP/FN breakdown guides improvement", 1, GRAY_LT),
    ]
    bullets(sl, takeaways, Inches(0.9), Inches(1.15), Inches(11.5), Inches(4.8))

    # Future work strip
    fw_box = sl.shapes.add_shape(1, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.42))
    fill(fw_box, CARD)
    txt(sl, "Future:  human-annotated GT  ·  dynamic categories  ·  LLM-as-Judge  ·  multi-project dashboard",
        Inches(0.8), Inches(6.16), Inches(11.8), Inches(0.3),
        size=13, color=GRAY_LT)

    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    potx = "/Users/sanderschulman/Developer/aiagentsproject/24-880 Presentation Template (1).potx"
    out  = "/Users/sanderschulman/Developer/aiagentsproject/EIA_Agent_System_Presentation.pptx"
    build(potx, out)
